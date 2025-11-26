import os
import re
import json
from typing import Dict, Any, List, Optional
from io import BytesIO

# Third-party libraries
from fastapi import FastAPI, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from openai import OpenAI
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import aiofiles
from dotenv import load_dotenv
# Load environment variables from .env file
load_dotenv()

# --- Configuration and Setup ---

# Initialize OpenAI client (API key is loaded from environment variable OPENAI_API_KEY)
try:
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
except Exception as e:
    print(f"OpenAI client initialization failed: {e}")
    client = None

# Define the themes and their corresponding file paths
THEMES = {
    "Marketing": "./upload/marketing.html",
    "Education": "./upload/education.html",
    "Portfolio": "./upload/portfolio.html",
    "Technology": "./upload/technology.html",
    "Business": "./upload/business.html",
}

# --- Pydantic Models ---

class PresentationRequest(BaseModel):
    """Input model for presentation generation endpoints."""
    topic: str = Field(..., description="The topic/title of the presentation.")
    theme: str = Field(..., description="The theme of the presentation (e.g., Marketing, Education).")
    num_slides: int = Field(10, ge=1, le=10, description="The number of slides to generate content for (max 10).")

class PresentationResponse(BaseModel):
    """Output model for the HTML generation endpoint."""
    html_content: str

# --- Utility Functions ---

def get_template_path(theme: str) -> str:
    theme_key = next((k for k in THEMES.keys() if k.lower() == theme.lower()), None)
    if not theme_key:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid theme: {theme}. Must be one of {list(THEMES.keys())}"
        )
    path = THEMES[theme_key]
    if not os.path.exists(path):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Template file not found for theme: {theme}"
        )
    return path

async def generate_content_with_openai(html_template: str, topic: str, theme: str, num_slides: int) -> str:
    if not client:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="OpenAI client is not initialized. Check API key configuration."
        )

    soup = BeautifulSoup(html_template, 'html.parser')
    marked_soup = BeautifulSoup(html_template, 'html.parser')
    marked_slides = marked_soup.find_all('div', class_='slide-container')[:num_slides]
    
    placeholder_map = {}
    
    for slide_index, slide in enumerate(marked_slides):
        for element in slide.find_all(text=True):
            if element.parent.name not in ['script', 'style', 'title'] and element.strip():
                if not re.match(r'^\s*(\d+ / \d+|← Previous|Next →|Explore Our Vision|Learn More|Schedule Meeting|Download Brochure|f|in|𝕏|Company Image|Years in Business|Team Members|CAGR|Total Addressable Market|By 2025|Revenue Growth|Valuation|Series C Funding|Customer Retention|Magic Number|Average Contract Value|Gross Margin|Average Experience|Team Size|Global Offices|Product Development|Beta Testing|Market Launch|Geographic Expansion|Partnership Development|Sales Growth|M&A Activity|Integration|Synergy Realization|Regulatory Compliance|Financial Audit|IPO Launch)\s*$', element.strip()):
                    placeholder_id = f"SLIDE_{slide_index+1}_TEXT_{len(placeholder_map) + 1}"
                    placeholder_map[placeholder_id] = element.strip()
                    element.replace_with(f"[[{placeholder_id}]]")

    prompt_template = f"""
    You are an expert presentation content generator. Your task is to rewrite the content for a presentation on the topic: "{topic}" with the theme "{theme}".
    
    The presentation structure is based on the following HTML. I have replaced all text content that needs to be rewritten with unique placeholder IDs in the format [[SLIDE_X_TEXT_Y]].
    
    Your goal is to provide new, relevant content for each placeholder ID.
    
    Rules:
    1. The new content MUST be based on the topic "{topic}" and the theme "{theme}".
    2. You MUST return a single JSON object with the structure: {{"new_content": {{"SLIDE_1_TEXT_1": "New content for slide 1", ...}}}}.
    3. The keys of the 'new_content' dictionary MUST be the exact placeholder IDs (e.g., "SLIDE_1_TEXT_1").
    4. The values of the 'new_content' dictionary MUST be the new text content for that placeholder.
    5. The new content should be concise and professional, suitable for a presentation slide.
    6. DO NOT include any introductory text, markdown, code blocks, or extra text outside the single JSON object.
    
    Here is the map of placeholder IDs and their original text (for context on the element's purpose):
    {json.dumps(placeholder_map, indent=2)}
    
    
    Return ONLY the JSON object.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are an expert presentation content generator. You return a single JSON object with the new content for the provided placeholder IDs."},
                {"role": "user", "content": prompt_template}
            ],
        )
        
        response_text = response.choices[0].message.content
        
        async with aiofiles.open("./openai_response.log", mode="w") as f:
            await f.write(response_text)
            
        json_match = re.search(r'```json\s*(\{.*\})\s*```', response_text, re.DOTALL)
        if json_match:
            json_content = json_match.group(1)
        else:
            json_match = re.search(r'(\{.*\})', response_text, re.DOTALL)
            if not json_match:
                raise HTTPException(
                    status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                    detail="OpenAI did not return a valid JSON object."
                )
            json_content = json_match.group(0)
            
        try:
            new_content_map = json.loads(json_content)["new_content"]
        except (json.JSONDecodeError, KeyError) as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to parse OpenAI response: {e}. Raw content saved to openai_response.log"
            )
            
        final_html = str(marked_soup)
        for placeholder_id, new_text in new_content_map.items():
            final_html = re.sub(re.escape(f"[[{placeholder_id}]]"), new_text, final_html)

        return final_html
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"OpenAI API call failed: {e}"
        )

def html_to_ppt(html_content: str, theme_path: str) -> BytesIO:
    """
    Converts the AI-generated HTML slides into a PowerPoint presentation,
    attempting to use a corresponding PPTX template for theme application.
    """
    # Check for a corresponding PPTX template file (e.g., business.html -> business.pptx)
    base_name = os.path.splitext(theme_path)[0]
    pptx_template_path = base_name + ".pptx"

    if os.path.exists(pptx_template_path):
        prs = Presentation(pptx_template_path)
    else:
        prs = Presentation()

    soup = BeautifulSoup(html_content, 'html.parser')
    slide_containers = soup.find_all('div', class_='slide-container')

    # Use a slide layout from the loaded presentation's master to inherit the theme.
    try:
        slide_layout = prs.slide_layouts[1]  # Title and Content
    except IndexError:
        slide_layout = prs.slide_layouts[6]  # Blank

    for slide_index, slide_div in enumerate(slide_containers):
        slide = prs.slides.add_slide(slide_layout)

        # Use placeholders if available to inherit theme styling
        try:
            title_shape = slide.shapes.title
        except AttributeError:
            title_shape = None

        try:
            # Try to use the body placeholder (index 1 is common for content)
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
        except (AttributeError, IndexError):
            # Fallback to a custom textbox
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()

        elements = slide_div.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li'])
        for i, elem in enumerate(elements):
            text = elem.get_text(strip=True)
            if not text:
                continue

            # Skip common non-content elements
            if re.match(r'^\s*(\d+ / \d+|← Previous|Next →|Explore Our Vision|Learn More|Schedule Meeting|Download Brochure|f|in|𝕏|Company Image|Years in Business|Team Members|CAGR|Total Addressable Market|By 2025|Revenue Growth|Valuation|Series C Funding|Customer Retention|Magic Number|Average Contract Value|Gross Margin|Average Experience|Team Size|Global Offices|Product Development|Beta Testing|Market Launch|Geographic Expansion|Partnership Development|Sales Growth|M&A Activity|Integration|Synergy Realization|Regulatory Compliance|Financial Audit|IPO Launch)\s*$', text):
                continue

            if title_shape and i == 0:
                title_shape.text = text
            else:
                # Add a new paragraph for the content
                if i == 0 and not title_shape:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    
                p.text = text
                # Apply basic formatting for text elements
                if elem.name.startswith('h'):
                    p.font.bold = True
                    p.level = 0
                    if elem.name == 'h1':
                        p.font.size = Pt(32)
                    elif elem.name == 'h2':
                        p.font.size = Pt(28)
                    else:
                        p.font.size = Pt(24)
                else:
                    p.level = 1 # Bullet point for p and li
                    p.font.size = Pt(20)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- FastAPI Application ---

app = FastAPI(
    title="AI Presentation Generator",
    description="A service to generate themed presentations in HTML and PowerPoint formats using OpenAI.",
    version="1.0.0"
)

# Add CORS middleware to allow all origins for development/testing
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    return {"message": "AI Presentation Generator API is running. See /docs for endpoints."}

from fastapi.responses import HTMLResponse

@app.post("/generate-html", response_class=HTMLResponse)
async def generate_html_presentation(request: PresentationRequest):
    """
    Generates an HTML presentation with AI-filled content based on the topic and theme.
    """
    theme_path = get_template_path(request.theme)
    
    async with aiofiles.open(theme_path, mode="r") as f:
        html_template = await f.read()
        
    final_html = await generate_content_with_openai(
        html_template, 
        request.topic, 
        request.theme, 
        request.num_slides
    )
    
    return HTMLResponse(content=final_html, status_code=200)

@app.post("/generate-ppt")
async def generate_ppt_presentation(request: PresentationRequest):
    """
    Generates a PowerPoint (.pptx) presentation with AI-filled content and returns it as a file download.
    """
    # 1. Generate the HTML content first
    theme_path = get_template_path(request.theme)
    async with aiofiles.open(theme_path, mode="r") as f:
        html_template = await f.read()
        
    final_html = await generate_content_with_openai(
        html_template, 
        request.topic, 
        request.theme, 
        request.num_slides
    )
    
    # 2. Convert HTML to PPTX
    ppt_io = html_to_ppt(final_html, theme_path)
    
    # 3. Return the PPTX file as a download
    from fastapi.responses import StreamingResponse
    
    filename = f"{request.theme}_{request.topic.replace(' ', '_')}.pptx"
    
    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )
