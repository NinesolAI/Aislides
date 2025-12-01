import os
import re
import json
from typing import Dict, Any, List, Optional
from io import BytesIO

# Third-party libraries
from fastapi import FastAPI, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from openai import OpenAI
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import aiofiles
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# --- Configuration and Setup ---

# Initialize OpenAI client (API key is loaded from environment variable OPENAI_API_KEY)
try:
    # client = OpenAI()
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
except Exception as e:
    # This is a placeholder for the actual error handling, as the client is initialized without args
    # because the API key is pre-configured in the environment.
    print(f"OpenAI client initialization failed: {e}")
    client = None

# Define the themes and their corresponding file paths
THEMES = {
    "Marketing": "./upload/marketing.html",
    "Education": "./upload/education.html",
    "Portfolio": "./upload/portfolio.html",
    "Technology": "./upload/technology.html",
    "Business": "./upload/business.html",
    "Businessv2": "./upload/businessv2.html",
}

# Define the PPTX template path
THEMESPPT = {
    "Technology": "./upload/technology.pptx",
    "Business": "./upload/businessv2.pptx",
}

# --- Pydantic Models ---

class PresentationRequest(BaseModel):
    """Input model for presentation generation endpoints."""
    topic: str = Field(..., description="The topic/title of the presentation.")
    theme: str = Field(..., description="The theme of the presentation (e.g., Marketing, Education).")
    num_slides: int = Field(10, ge=1, le=10, description="The number of slides to generate content for (max 10).")

class PresentationResponse(BaseModel):
    """Output model for the HTML generation endpoint."""
    html_content: str

# --- Utility Functions ---

def get_template_path(theme: str) -> str:
    theme_key = next((k for k in THEMES.keys() if k.lower() == theme.lower()), None)
    if not theme_key:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid theme: {theme}. Must be one of {list(THEMES.keys())}"
        )
    path = THEMES[theme_key]
    if not os.path.exists(path):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Template file not found for theme: {theme}"
        )
    return path

def get_pptx_template_path(theme: str) -> str:
    theme_key = next((k for k in THEMESPPT.keys() if k.lower() == theme.lower()), None)
    if not theme_key:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid PPTX theme: {theme}. Must be one of {list(THEMESPPT.keys())}"
        )
    path = THEMESPPT[theme_key]
    if not os.path.exists(path):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"PPTX template file not found for theme: {theme}"
        )
    return path

def pptx_text_processor(
    pptx_path: str, 
    new_content_map: Optional[Dict[str, str]] = None
) -> tuple[Dict[str, str], BytesIO]:
    """
    Extracts text from a PPTX file and prepares a placeholder map, 
    or replaces text with new content if a map is provided.
    
    Returns: (placeholder_map, pptx_io)
    - placeholder_map: A dictionary of {placeholder_id: original_text}
    - pptx_io: A BytesIO object of the modified PPTX (only if new_content_map is provided)
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX template not found at {pptx_path}")

    prs = Presentation(pptx_path)
    placeholder_map = {}
    
    # We will use a list of tuples to store references to the text runs
    # (slide_index, shape_index, paragraph_index, run_index)
    text_run_references = []
    
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            for para_index, paragraph in enumerate(text_frame.paragraphs):
                # We will process text at the run level for more granular replacement
                for run_index, run in enumerate(paragraph.runs):
                    text = run.text.strip()
                    if text:
                        # Create a unique ID for the placeholder
                        placeholder_id = f"S{slide_index+1}_SH{shape_index+1}_P{para_index+1}_R{run_index+1}"
                        
                        if new_content_map is None:
                            # Extraction mode: collect original text
                            placeholder_map[placeholder_id] = run.text
                            text_run_references.append((slide_index, shape_index, para_index, run_index))
                        else:
                            # Replacement mode: replace text if a new value is provided
                            new_text = new_content_map.get(placeholder_id)
                            if new_text is not None:
                                run.text = new_text
                                
    if new_content_map is None:
        # Return the map for content generation
        return placeholder_map, None
    else:
        # Save the modified presentation to a BytesIO object
        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        return placeholder_map, ppt_io

# The existing HTML utility functions (extract_placeholder_text, generate_content_with_openai, html_to_ppt) 
# are assumed to be correct and will be kept.

async def generate_content_with_openai(template_content: str, topic: str, theme: str, num_slides: int, is_pptx: bool = False) -> str:
    """
    Calls the OpenAI API to generate new content based on the template structure.
    Returns the full content with the new content (HTML string or JSON string for PPTX).
    """
    if not client:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="OpenAI client is not initialized. Check API key configuration."
        )

    if is_pptx:
        # PPTX logic: template_content is the placeholder map JSON string
        placeholder_map = json.loads(template_content)
        
        prompt_template = f"""
        You are an expert presentation content generator. Your task is to rewrite the content for a presentation on the topic: "{topic}" with the theme "{theme}".
        
        The presentation structure is based on the following text placeholders extracted from a PowerPoint file.
        
        Your goal is to provide new, relevant content for each placeholder ID.
        
        Rules:
        1. The new content MUST be based on the topic "{topic}" and the theme "{theme}".
        2. You MUST return a single JSON object with the structure: {{"new_content": {{"S1_SH1_P1_R1": "New content for slide 1", ...}}}}.
        3. The keys of the 'new_content' dictionary MUST be the exact placeholder IDs (e.g., "S1_SH1_P1_R1").
        4. The values of the 'new_content' dictionary MUST be the new text content for that placeholder.
        5. The new content should be concise and professional, suitable for a presentation slide.
        6. DO NOT include any introductory text, markdown, code blocks, or extra text outside the single JSON object.
        7. Only generate content for the first {num_slides} slides.
        
        Here is the map of placeholder IDs and their original text (for context on the element's purpose):
        {json.dumps(placeholder_map, indent=2)}
        
        
        Return ONLY the JSON object.
        """
    else:
        # HTML logic (re-using the existing logic)
        # ... (The existing HTML logic for prompt generation goes here)
        
        # 1. Extract the HTML structure for the first `num_slides`
        soup = BeautifulSoup(template_content, 'html.parser')
        
        # Create a copy of the soup for marking
        marked_soup = BeautifulSoup(template_content, 'html.parser')
        marked_slides = marked_soup.find_all('div', class_='slide-container')[:num_slides]
        
        placeholder_map = {}
        
        for slide_index, slide in enumerate(marked_slides):
            # Find all text nodes that are not inside script or style tags
            for element in slide.find_all(text=True):
                if element.parent.name not in ['script', 'style', 'title'] and element.strip():
                    # Check if the text is not a navigation element (like '← Previous')
                    if not re.match(r'^\s*(\d+ / \d+|← Previous|Next →|Explore Our Vision|Learn More|Schedule Meeting|Download Brochure|f|in|𝕏|Company Image|Years in Business|Team Members|CAGR|Total Addressable Market|By 2025|Revenue Growth|Valuation|Series C Funding|Customer Retention|Magic Number|Average Contract Value|Gross Margin|Average Experience|Team Size|Global Offices|Product Development|Beta Testing|Market Launch|Geographic Expansion|Partnership Development|Sales Growth|M&A Activity|Integration|Synergy Realization|Regulatory Compliance|Financial Audit|IPO Launch)\s*$', element.strip()):
                        # Create a unique ID for the placeholder
                        placeholder_id = f"SLIDE_{slide_index+1}_TEXT_{len(placeholder_map) + 1}"
                        placeholder_map[placeholder_id] = element.strip()
                        
                        # Replace the text node with the unique ID marker
                        element.replace_with(f"[[{placeholder_id}]]")

        # The prompt will ask the model to fill in the content for the marked IDs.
        prompt_template = f"""
        You are an expert presentation content generator. Your task is to rewrite the content for a presentation on the topic: "{topic}" with the theme "{theme}".
        
        The presentation structure is based on the following HTML. I have replaced all text content that needs to be rewritten with unique placeholder IDs in the format [[SLIDE_X_TEXT_Y]].
        
        Your goal is to provide new, relevant content for each placeholder ID.
        
        Rules:
        1. The new content MUST be based on the topic "{topic}" and the theme "{theme}".
        2. You MUST return a single JSON object with the structure: {{"new_content": {{"SLIDE_1_TEXT_1": "New content for slide 1", ...}}}}.
        3. The keys of the 'new_content' dictionary MUST be the exact placeholder IDs (e.g., "SLIDE_1_TEXT_1").
        4. The values of the 'new_content' dictionary MUST be the new text content for that placeholder.
        5. The new content should be concise and professional, suitable for a presentation slide.
        6. DO NOT include any introductory text, markdown, code blocks, or extra text outside the single JSON object.
        
        Here is the map of placeholder IDs and their original text (for context on the element's purpose):
        {json.dumps(placeholder_map, indent=2)}
        
        
        Return ONLY the JSON object.
        """
        
    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini", # Using a fast model for content generation
            messages=[
                {"role": "system", "content": "You are an expert presentation content generator. You return a single JSON object with the new content for the provided placeholder IDs."},
                {"role": "user", "content": prompt_template}
            ],
            # response_format={"type": "json_object"} # This is not supported by the current client version
        )
        
        # Parse the JSON response
        response_text = response.choices[0].message.content
        
        # Write the raw response to a file for debugging
        async with aiofiles.open("./openai_response.log", mode="w") as f:
            await f.write(response_text)
            
        # Attempt to extract JSON from the response text
        json_match = re.search(r'```json\s*(\{.*\})\s*```', response_text, re.DOTALL)
        if json_match:
            json_content = json_match.group(1)
        else:
            # Fallback for non-markdown JSON
            json_match = re.search(r'(\{.*\})', response_text, re.DOTALL)
            if not json_match:
                raise HTTPException(
                    status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                    detail="OpenAI did not return a valid JSON object."
                )
            json_content = json_match.group(0)
            
        try:
            new_content_map = json.loads(json_content)["new_content"]
        except (json.JSONDecodeError, KeyError) as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to parse OpenAI response: {e}. Raw content saved to openai_response.log"
            )
            
        if is_pptx:
            # For PPTX, we return the new content map as a JSON string
            return json.dumps(new_content_map)
        else:
            # For HTML, we replace the markers and return the final HTML
            final_html = str(marked_soup)
            for placeholder_id, new_text in new_content_map.items():
                # Escape special characters in the new text for safe regex replacement
                # We use re.escape on the placeholder ID to ensure it's treated literally
                # We use re.sub to replace the marker with the new content
                final_html = re.sub(re.escape(f"[[{placeholder_id}]]"), new_text, final_html)

            return final_html
            
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred during content generation: {e}"
        )

# The existing html_to_ppt function is assumed to be correct and will be kept.
def html_to_ppt(html_content: str) -> BytesIO:
    """
    Converts the AI-generated HTML slides into a PowerPoint presentation with
    basic styling: headings bolded, bullet points, and font sizes.
    """
    prs = Presentation()
    soup = BeautifulSoup(html_content, 'html.parser')
    slide_containers = soup.find_all('div', class_='slide-container')

    # Use a blank slide layout for flexibility
    blank_slide_layout = prs.slide_layouts[6]

    for slide_index, slide_div in enumerate(slide_containers):
        slide = prs.slides.add_slide(blank_slide_layout)

        # Define a text box for the content
        left = top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Extract headings and paragraphs
        elements = slide_div.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li'])
        for i, elem in enumerate(elements):
            text = elem.get_text(strip=True)
            if not text:
                continue

            # Determine font size based on heading level
            if elem.name == 'h1':
                font_size = Pt(36)
                bold = True
            elif elem.name == 'h2':
                font_size = Pt(32)
                bold = True
            elif elem.name == 'h3':
                font_size = Pt(28)
                bold = True
            elif elem.name == 'h4':
                font_size = Pt(24)
                bold = True
            else:  # h5, h6, p, li
                font_size = Pt(20)
                bold = False

            # First element in the text frame uses the first paragraph
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            p.font.size = font_size
            p.font.bold = bold
            p.level = 0 if elem.name.startswith('h') else 1  # indent bullet points for non-headings
            p.font.color.rgb = RGBColor(0, 0, 0)  # black text by default

    # Save presentation to BytesIO
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- FastAPI Application ---

app = FastAPI(
    title="AI Presentation Generator",
    description="A service to generate themed presentations in HTML and PowerPoint formats using OpenAI.",
    version="1.0.0"
)

# Add CORS middleware to allow all origins for development/testing
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    return {"message": "AI Presentation Generator API is running. See /docs for endpoints."}

from fastapi.responses import HTMLResponse, StreamingResponse

@app.post("/generate-html", response_class=HTMLResponse)
async def generate_html_presentation(request: PresentationRequest):
    """
    Generates an HTML presentation with AI-filled content based on the topic and theme.
    """
    theme_path = get_template_path(request.theme)
    
    # Read the HTML template asynchronously
    async with aiofiles.open(theme_path, mode="r") as f:
        html_template = await f.read()
        
    # Generate content and replace placeholders
    final_html = await generate_content_with_openai(
        html_template, 
        request.topic, 
        request.theme, 
        request.num_slides,
        is_pptx=False
    )
    
    return HTMLResponse(content=final_html, status_code=200)

@app.post("/generate-ppt")
async def generate_ppt_presentation(request: PresentationRequest):
    """
    Generates a PowerPoint (.pptx) presentation with AI-filled content and returns it as a file download.
    This endpoint converts an AI-generated HTML presentation to PPTX.
    """
    # 1. Generate the HTML content first
    theme_path = get_template_path(request.theme)
    async with aiofiles.open(theme_path, mode="r") as f:
        html_template = await f.read()
        
    final_html = await generate_content_with_openai(
        html_template, 
        request.topic, 
        request.theme, 
        request.num_slides,
        is_pptx=False
    )
    
    # 2. Convert HTML to PPTX
    ppt_io = html_to_ppt(final_html)
    
    # 3. Return the PPTX file as a download
    filename = f"{request.theme}_{request.topic.replace(' ', '_')}.pptx"
    
    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/generate-pptx-from-template")
async def generate_pptx_from_template(request: PresentationRequest):
    """
    Generates a PowerPoint (.pptx) presentation by filling content into the uploaded technology.pptx template.
    """
    # 1. Extract text and create placeholder map from the template
    try:
        pptx_path = get_pptx_template_path(request.theme)
        placeholder_map, _ = pptx_text_processor(pptx_path, new_content_map=None)
    except FileNotFoundError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e)
        )
    
    # Filter the map to only include placeholders from the requested number of slides
    filtered_map = {
        k: v for k, v in placeholder_map.items() 
        if int(k.split('_')[0].replace('S', '')) <= request.num_slides
    }
    
    # 2. Generate new content using OpenAI
    # The template_content for PPTX is the JSON string of the placeholder map
    new_content_json_str = await generate_content_with_openai(
        json.dumps(filtered_map), 
        request.topic, 
        request.theme, 
        request.num_slides,
        is_pptx=True
    )
    
    new_content_map = json.loads(new_content_json_str)
    
    # 3. Replace text in the PPTX template
    try:
        _, ppt_io = pptx_text_processor(pptx_path, new_content_map=new_content_map)
    except FileNotFoundError as e:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(e)
        )
        
    # 4. Return the modified PPTX file as a download
    filename = f"{request.theme}_{request.topic.replace(' ', '_')}_template_filled.pptx"
    
    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# To run the application: uvicorn app:app --host 0.0.0.0 --port 8000